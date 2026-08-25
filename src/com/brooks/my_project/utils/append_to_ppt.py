from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import time


def append_pages_to_ppt():
    """
    打开现有PPT，追加内容丰富的幻灯片以增大文件体积。
    基于：岚图汽车党委关于开展深入贯彻中央八项规定精神学习教育的实施进展
    """
    # ===================== 文件路径配置 =====================
    input_file = r"D:\ActionSoft\workspace\code\MY_PYTHON\src\com\brooks\my_project\岚图汽车党委关于开展深入贯彻中央八项规定精神学习教育的实施进展.pptx"
    output_file = r"D:\ActionSoft\workspace\code\MY_PYTHON\src\com\brooks\my_project\岚图汽车党委关于开展深入贯彻中央八项规定精神学习教育的实施进展_补全.pptx"
    target_size_mb = 80
    # ========================================================

    print("=" * 50)
    print("📂 正在打开现有文件...")
    print(f"   文件路径：{input_file}")

    if not os.path.exists(input_file):
        print(f"❌ 错误：找不到文件！")
        return

    prs = Presentation(input_file)

    # 获取当前文件大小
    current_size = os.path.getsize(input_file) / (1024 * 1024)
    print(f"📊 当前文件大小：{current_size:.2f} MB")

    if current_size >= target_size_mb:
        print("✅ 文件已经达到目标大小，无需追加。")
        return

    need_size = target_size_mb - current_size
    print(f"🎯 需要增加：{need_size:.2f} MB")

    # 估算页数：每页约12KB (0.012 MB)
    estimated_pages = int(need_size / 0.012) + 200
    print(f"📄 估算需要追加约 {estimated_pages} 页")
    print(f"⏱️  预计耗时约 {estimated_pages / 50:.0f} 秒")

    # ---------- 构建丰富的内容（基于新PPT） ----------
    base_text = (
        "岚图汽车党委关于深入贯彻中央八项规定精神学习教育的实施进展。"
        "工作背景：3月19日东风公司党委召开研究部署会，全面启动相关工作。"
        "3月24日岚图汽车召开党委会，研究审议工作方案并进行动员部署。"
        "3月25日至27日各党支部以三会一课主题党日方式进行启动部署。"
        "4月8日召开学习教育工作推进会，各党工团组织负责人参加。"
        "坚持学习在先，思想引领：深入学习中央八项规定精神，提高政治站位。"
        "坚持深查问题，以案为鉴：对照规定要求，深入查摆存在的问题。"
        "坚持以改促治，边查边改：针对问题制定整改措施，推动制度完善。"
        "坚持开门教育，群众参与：广泛听取群众意见，接受群众监督。"
        "坚持常态推进，深查细悟：建立长效机制，持续深化学习教育。"
        "提请经管会决策：对开展学习教育提出具体要求。"
    )
    content_lines = [f"{i+1:02d}. {base_text}" for i in range(30)]

    # ---------- 安全获取空白版式 ----------
    # 尝试索引6（常见空白版式），如果不存在则使用最后一个版式
    try:
        slide_layout = prs.slide_layouts[6]
    except IndexError:
        # 如果没有索引6，使用最后一个版式（通常也是空白或内容版式）
        slide_layout = prs.slide_layouts[-1]
    print(f"📐 使用版式索引：{prs.slide_layouts.index(slide_layout)}")

    print("⏳ 开始追加页面...")
    start_time = time.time()

    for i in range(estimated_pages):
        slide = prs.slides.add_slide(slide_layout)
        # 添加文本框（覆盖几乎整页）
        shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(6.5))
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.text = "\n".join(content_lines)
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = "微软雅黑"
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
            paragraph.alignment = PP_ALIGN.LEFT

        if (i + 1) % 100 == 0 or (i + 1) == estimated_pages:
            elapsed = time.time() - start_time
            speed = (i + 1) / elapsed if elapsed > 0 else 0
            print(f"  已添加 {i+1}/{estimated_pages} 页，速度 {speed:.1f} 页/秒，已用时 {elapsed:.0f} 秒")

    # 保存文件
    print("💾 正在保存PPT文件...")
    prs.save(output_file)

    new_size = os.path.getsize(output_file) / (1024 * 1024)
    total_time = time.time() - start_time
    print("=" * 50)
    print(f"✅ 追加完成！")
    print(f"📁 新文件：{os.path.abspath(output_file)}")
    print(f"📊 新文件大小：{new_size:.2f} MB")
    print(f"⏱️  追加耗时：{total_time:.1f} 秒")

    if new_size < target_size_mb:
        print(f"⚠️ 还差 {target_size_mb - new_size:.2f} MB")
        print(f"💡 建议：再运行一次脚本，或手动插入几张图片")
    else:
        print("🎉 已达成目标 60MB！")
    print("=" * 50)


if __name__ == "__main__":
    append_pages_to_ppt()