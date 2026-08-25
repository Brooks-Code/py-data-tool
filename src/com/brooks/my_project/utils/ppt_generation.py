from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os
import time
from typing import List, Optional, Tuple


# ==================== 配置类 ====================
class PPTConfig:
    """PPT生成配置"""
    # 颜色
    RED = RGBColor(192, 0, 0)
    DARK_RED = RGBColor(128, 0, 0)
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)
    GRAY = RGBColor(100, 100, 100)
    LIGHT_GRAY = RGBColor(240, 240, 240)

    # 尺寸
    SLIDE_WIDTH = 13.333
    SLIDE_HEIGHT = 7.5
    FONT_NAME = "微软雅黑"

    # 重复页配置
    REPEAT_PAGES = 5000  # 可通过外部调整
    TARGET_SIZE_MB = 60


class PPTBuilder:
    """PPT构建器 - 封装所有操作"""

    def __init__(self, filename: str = "审议增强自信决战决胜专项思想政治工作方案.pptx"):
        self.filename = filename
        self.prs = Presentation()
        self.prs.slide_width = Inches(PPTConfig.SLIDE_WIDTH)
        self.prs.slide_height = Inches(PPTConfig.SLIDE_HEIGHT)
        self.page_count = 0

        # 颜色快捷引用
        self.C = PPTConfig

    def add_text_box(self, slide, text: str, left: float, top: float,
                     width: float, height: float, font_size: int = 18,
                     font_color=PPTConfig.BLACK, bold: bool = False,
                     align=PP_ALIGN.LEFT, fill_color=None) -> None:
        """添加文本框"""
        shape = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        text_frame = shape.text_frame
        text_frame.text = str(text)
        text_frame.word_wrap = True

        for paragraph in text_frame.paragraphs:
            paragraph.alignment = align
            paragraph.font.name = PPTConfig.FONT_NAME
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = font_color
            paragraph.font.bold = bold

        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()

    def add_red_bar(self, slide, x: float, y: float, width: float, height: float) -> None:
        """添加红色装饰条"""
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = PPTConfig.RED
        rect.line.fill.background()

    def add_title_slide(self, title: str, subtitle: str = "", sub_info: str = "") -> None:
        """添加标题幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.page_count += 1

        # 顶部红色条
        self.add_red_bar(slide, 0, 0, PPTConfig.SLIDE_WIDTH, 0.5)

        # 主标题
        self.add_text_box(slide, title, 1, 1.8, 11.333, 1.5,
                          font_size=40, font_color=PPTConfig.RED,
                          bold=True, align=PP_ALIGN.CENTER)

        # 副标题
        if subtitle:
            self.add_text_box(slide, subtitle, 1, 3.2, 11.333, 0.8,
                              font_size=24, font_color=PPTConfig.BLACK,
                              align=PP_ALIGN.CENTER)

        # 底部信息
        if sub_info:
            self.add_text_box(slide, sub_info, 1, 5.8, 11.333, 0.6,
                              font_size=16, font_color=PPTConfig.GRAY,
                              align=PP_ALIGN.CENTER)

        # 底部红色条
        self.add_red_bar(slide, 0, 7.0, PPTConfig.SLIDE_WIDTH, 0.5)

    def add_content_slide(self, title: str, content_lines: List[str],
                          sub_title: str = "") -> None:
        """添加内容幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.page_count += 1

        # 左侧红色条
        self.add_red_bar(slide, 0, 0.5, 0.15, 6.5)

        # 标题
        self.add_text_box(slide, title, 0.8, 0.4, 11, 0.8,
                          font_size=28, font_color=PPTConfig.DARK_RED, bold=True)

        # 副标题
        if sub_title:
            self.add_text_box(slide, sub_title, 0.8, 1.1, 11, 0.5,
                              font_size=16, font_color=PPTConfig.GRAY)

        # 内容
        y_pos = 1.6 if sub_title else 1.4
        for line in content_lines:
            if y_pos > 6.2:
                break

            stripped = line.strip()
            if not stripped:
                y_pos += 0.25
                continue

            # 判断缩进级别
            if stripped.startswith(("-", "•", "", "")):
                self.add_text_box(slide, line, 1.0, y_pos, 11.5, 0.45,
                                  font_size=18, font_color=PPTConfig.BLACK, bold=True)
            elif stripped.startswith(("  ", "    ")):
                self.add_text_box(slide, line, 1.4, y_pos, 11.0, 0.4,
                                  font_size=16, font_color=PPTConfig.BLACK)
            else:
                self.add_text_box(slide, line, 1.0, y_pos, 11.0, 0.4,
                                  font_size=16, font_color=PPTConfig.BLACK)
            y_pos += 0.45

    def add_table_slide(self, title: str, table_data: List[Tuple[str, ...]]) -> None:
        """添加表格幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.page_count += 1

        # 左侧红色条
        self.add_red_bar(slide, 0, 0.5, 0.15, 6.5)

        # 标题
        self.add_text_box(slide, title, 0.8, 0.4, 11, 0.8,
                          font_size=24, font_color=PPTConfig.DARK_RED, bold=True)

        if not table_data:
            return

        # 列配置
        x_positions = [0.8, 1.8, 5.3, 8.3]
        widths = [1.0, 3.5, 3.0, 2.5]

        # 表头
        headers = table_data[0]
        for i, (x, w) in enumerate(zip(x_positions, widths)):
            self.add_text_box(
                slide, headers[i] if i < len(headers) else "",
                x, 1.4, w, 0.5,
                font_size=16, font_color=PPTConfig.WHITE, bold=True,
                align=PP_ALIGN.CENTER, fill_color=PPTConfig.RED
            )

        # 数据行
        y_pos = 1.9
        for idx, row in enumerate(table_data[1:], 1):
            bg_color = PPTConfig.LIGHT_GRAY if idx % 2 == 0 else PPTConfig.WHITE
            for i, (x, w) in enumerate(zip(x_positions, widths)):
                align = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT
                self.add_text_box(
                    slide, row[i] if i < len(row) else "",
                    x, y_pos, w, 0.45,
                    font_size=14, font_color=PPTConfig.BLACK,
                    align=align, fill_color=bg_color
                )
            y_pos += 0.48

    def add_repeated_slides(self, count: int,
                            title_template: str = "附加页 {}",
                            content: List[str] = None) -> None:
        """
        批量添加重复幻灯片（优化版）

        Args:
            count: 重复页数
            title_template: 标题模板，{} 会被替换为页码
            content: 内容列表
        """
        if content is None:
            content = [
                "本页为重复内容，用于增加文件体积。",
                "实际内容可以根据需要修改。",
                "此方法可快速增大PPT文件大小。"
            ]

        print(f"⏳ 开始添加 {count} 页重复内容...")
        start_time = time.time()

        # 批量生成，每100页输出一次进度
        batch_size = 100
        for i in range(1, count + 1):
            self.add_content_slide(
                title_template.format(i),
                content
            )

            # 进度提示
            if i % batch_size == 0 or i == count:
                elapsed = time.time() - start_time
                speed = i / elapsed if elapsed > 0 else 0
                print(f"  已添加 {i}/{count} 页，速度: {speed:.1f} 页/秒")

        print(f"✅ 重复页添加完成，耗时: {time.time() - start_time:.1f} 秒")

    def build(self) -> str:
        """构建并保存PPT"""
        print("⏳ 开始生成思想政治工作方案PPT...")
        total_start = time.time()

        # ====== 第1页：封面 ======
        self.add_title_slide(
            "审议",
            '"增强自信、决战决胜"\n专项思想政治工作方案',
            "汇报专业：党群工作部 ｜ 汇报人：张宇"
        )
        print("✅ 第1页：封面")

        # ====== 第2页：目录 ======
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.page_count += 1
        self.add_red_bar(slide, 0, 0.5, 0.15, 6.5)

        self.add_text_box(slide, "目录", 0.8, 0.4, 3, 0.8,
                          font_size=32, font_color=PPTConfig.DARK_RED, bold=True)

        items = ["一、议题背景", "二、议题内容", "三、提请审议"]
        y_pos = 1.8
        for idx, item in enumerate(items, 1):
            self.add_text_box(slide, f"0{idx}", 0.8, y_pos, 0.8, 0.8,
                              font_size=36, font_color=PPTConfig.RED, bold=True)
            self.add_text_box(slide, item, 2.0, y_pos + 0.05, 8, 0.7,
                              font_size=28, font_color=PPTConfig.BLACK)

            # 分隔线
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(2.0), Inches(y_pos + 0.8),
                Inches(8), Inches(0.02)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = PPTConfig.LIGHT_GRAY
            line.line.fill.background()
            y_pos += 1.2
        print("✅ 第2页：目录")

        # ====== 第3-10页：主要内容 ======
        content_pages = [
            ("一、议题背景", [
                "", '以增强"技术自信、产品自信、品牌自信、价值自信"为主线，',
                "推动思想淬炼与战略执行深度融合。", "", "• 群团暖心", "• 组织保障", "• 宣传动员"
            ]),
            ("二、议题内容——宣传动员", [
                "01  让榜样站出来说话",
                "    策划'青年说''梦想，这么干'等宣传主题，",
                "    用获得荣誉的团队和个人走到台前，用一线事迹感染人。",
                "    同步汇编《奋斗之星故事集》，把典型树在员工身边。", "",
                "02  用爆款内容打透认知",
                "    集中火力推'泰山X8天梯爬坡'全媒体视频传播。",
                "    策划FREE下线五周年专题，把产品实力、技术突破",
                "    直观呈现到员工眼前。"
            ]),
            ("二、议题内容——宣传动员（续）", [
                "03  用快讯抢占注意力",
                "    开设'We are 岚波one'专栏，每周短平快推送",
                "    形势任务、一线动态，不让杂音先入为主。", "",
                "04  主题宣传片点燃斗志",
                "    高标准制作七一主题宣传片《使命》，选取8-10名",
                "    一线普通党员的故事，在七一表彰大会上首发。", "",
                "05  企业文化节系列活动",
                "    系统整理话术体系，组织岚图好声音大赛，",
                "    在活动中升温感情、凝聚认同。"
            ]),
            ("二、议题内容——组织保障", [
                "01  学习贯彻习近平党建思想",
                "    落实中央要求，组织学习习近平党建思想，",
                "    梳理营销领域党组织设置情况、党员分布情况。", "",
                "02  建强基层党组织",
                "    动态优化组织设置，加大党员培养和教育力度，",
                "    充分发挥党支部战斗堡垒作用和党员先锋模范作用。", "",
                "03  开展一次大讨论",
                "    深入开展'增强四个自信、决战决胜下半年'",
                "    支部大讨论，聚焦增强四个自信，剖析年度目标。"
            ]),
            ("二、议题内容——群团暖心", [
                "• 新入职大学生快速融入",
                "    入职第一课专题讲授岚图文化，",
                "    广泛开展文体协会纳新，帮助新员工加速融入。", "",
                "• 高温送清凉慰问",
                "    针对三个一线人员，开展'战高温、保高产'专项慰问。", "",
                "• 员工关怀升级",
                "    端午福利标准升级至500元/人。",
                "    升级生日观影福利项目，标准提至600元/人。"
            ]),
            ("二、议题内容——群团暖心（续）", [
                "• 家企连心暑期关爱行动",
                "    开办'岚精灵'暑期托管班，解决员工子女看护难题。",
                "    组织'员工家属开放日'，邀请家人走进公司。", "",
                "• EAP心理关爱专项活动",
                "    针对重点专业、关键人群，开展EAP心理关爱专项活动，",
                "    压力疏导、心理咨询、团体辅导。", "",
                "• 针对制造领域特殊诉求，做好针对性服务。"
            ]),
            ("三、提请审议", [
                "", "请党委会审议", "", '"增强自信、决战决胜"', "专项思想政治工作方案"
            ]),
        ]

        for idx, (title, content) in enumerate(content_pages, start=3):
            self.add_content_slide(title, content)
            print(f"✅ 第{idx}页：{title}")

        # ====== 第11页：行动排期表 ======
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.page_count += 1
        self.add_red_bar(slide, 0, 0.5, 0.15, 6.5)

        self.add_text_box(slide, "二、议题内容——行动排期表", 0.8, 0.4, 11, 0.8,
                          font_size=24, font_color=PPTConfig.DARK_RED, bold=True)

        table_data = [
            ("序号", "行动事项", "责任部门", "时间节点"),
            ("1", "宣传动员系列工作", "党群工作部", "7月-9月"),
            ("2", "组织保障系列工作", "党群工作部/人力", "7月-12月"),
            ("3", "群团暖心系列工作", "党群工作部", "7月-10月"),
            ("4", "企业文化节活动", "党群工作部", "8月"),
            ("5", "七一主题宣传片发布", "党群工作部", "7月1日"),
        ]
        self.add_table_slide("", table_data)  # 标题已在前面添加
        print("✅ 第11页：行动排期表")

        # ====== 感谢页 ======
        self.add_title_slide("感谢倾听", "Thanks", "党群工作部 ｜ 张宇")
        print("✅ 第12页：感谢页")

        # ====== 批量添加重复页 ======
        self.add_repeated_slides(PPTConfig.REPEAT_PAGES)

        # ====== 保存 ======
        print("⏳ 正在保存PPT文件...")
        self.prs.save(self.filename)

        total_time = time.time() - total_start
        abs_path = os.path.abspath(self.filename)
        file_size = os.path.getsize(self.filename) / (1024 * 1024)

        print(f"\n{'=' * 50}")
        print(f"🎉 PPT生成完成！")
        print(f"📁 文件保存至：{abs_path}")
        print(f"📄 总页数：{self.page_count} 页")
        print(f"📊 文件大小：{file_size:.2f} MB")
        print(f"⏱️  总耗时：{total_time:.1f} 秒")

        # 大小评估
        if file_size < PPTConfig.TARGET_SIZE_MB * 0.9:
            needed_pages = int(PPTConfig.REPEAT_PAGES * (PPTConfig.TARGET_SIZE_MB / file_size))
            print(f"\n⚠️ 当前大小为 {file_size:.1f} MB，目标 {PPTConfig.TARGET_SIZE_MB} MB")
            print(f"💡 建议将 REPEAT_PAGES 从 {PPTConfig.REPEAT_PAGES} 调整为 {needed_pages}")
        elif file_size > PPTConfig.TARGET_SIZE_MB * 1.1:
            needed_pages = int(PPTConfig.REPEAT_PAGES * (PPTConfig.TARGET_SIZE_MB / file_size))
            print(f"\n⚠️ 当前大小为 {file_size:.1f} MB，超过目标 {PPTConfig.TARGET_SIZE_MB} MB")
            print(f"💡 建议将 REPEAT_PAGES 从 {PPTConfig.REPEAT_PAGES} 调整为 {needed_pages}")
        else:
            print(f"\n✅ 文件大小 {file_size:.1f} MB，已达到目标！")
        print(f"{'=' * 50}")

        return abs_path


# ==================== 主函数 ====================
def create_ideological_work_ppt(
        filename: str = "审议增强自信决战决胜专项思想政治工作方案.pptx",
        repeat_pages: int = 5000
) -> str:
    """
    生成思想政治工作方案PPT

    Args:
        filename: 输出文件名
        repeat_pages: 重复页数量（用于控制文件大小）
    """
    # 配置重复页数
    PPTConfig.REPEAT_PAGES = repeat_pages

    # 构建PPT
    builder = PPTBuilder(filename)
    return builder.build()


if __name__ == "__main__":
    # 你可以在这里调整 repeat_pages 来控制文件大小
    # 根据运行结果调整：每1000页约增加 2-5 MB
    create_ideological_work_ppt(repeat_pages=40000)